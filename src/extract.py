import hashlib
import io
import re
from dataclasses import dataclass
from typing import Optional, Tuple

import pdfplumber
from docx import Document
from pptx import Presentation


def stable_id_from_bytes(content: bytes) -> str:
    # deterministic id to avoid duplicates
    return hashlib.sha256(content).hexdigest()[:16]


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    parts.append(t)
    return "\n".join(parts).strip()


def extract_text_from_pdf(file_bytes: bytes) -> str:
    text = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            t = page.extract_text() or ""
            t = t.strip()
            if t:
                text.append(t)
    return "\n".join(text).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    return "\n".join(parts).strip()


def extract_text_generic(filename: str, file_bytes: bytes) -> Tuple[str, str]:
    """Return (doc_type, text)."""
    fn = filename.lower()
    if fn.endswith(".pptx"):
        return "pptx", extract_text_from_pptx(file_bytes)
    if fn.endswith(".pdf"):
        return "pdf", extract_text_from_pdf(file_bytes)
    if fn.endswith(".docx"):
        return "docx", extract_text_from_docx(file_bytes)
    if fn.endswith(".txt"):
        return "txt", file_bytes.decode("utf-8", errors="replace")
    raise ValueError("Unsupported file type. Use .pptx, .pdf, .docx or .txt")
